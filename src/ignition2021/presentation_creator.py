from pptx import *
from pptx.util import *
from pptx.enum.text import MSO_AUTO_SIZE
import os
from .presentation_reader import PresentationReader


class PresentationCreator:
    def __init__(self, presentation, image_dir):
        self.filename = presentation
        self.image_dir = f"ppts/{image_dir}"
        self.prs = Presentation(f"ppts/{presentation}")
        self.prs_r = PresentationReader(f"ppts/{presentation}", complexity=1, debug=True)
        self.margin = 10


    def ReadPresentation(self):
        self.prs_r.set_download_dir(self.image_dir)
        self.prs_r.set_image_quality("regular")
        self.prs_r.find_images()
        

    def maximalRectangle(self, matrix):
        if not matrix:
            return 0
        t = [n for n in matrix[0]]+[0]
        a, b, c = 0, 0, 0
        ans = 0
        for k in range(self.margin, len(matrix)):
            for j in range(len(matrix[0])):
                if matrix[k][j] == 0:
                    t[j] = 0
                else:
                    t[j] += int(matrix[k][j])
            mx = 0
            lft, rit, bot = 0, 0, 0
            stack = []
            i = 0
            while i < len(t):
                if len(stack) == 0 or t[stack[-1]] <= t[i]:
                    stack.append(i)
                    i += 1
                else:
                    index = stack.pop()
                    h = t[index]
                    w = i
                    if stack:
                        w = i - stack[-1] - 1
                    if w * h > mx:
                        mx = max(mx, w * h)
                        lft = stack[-1] if stack else 0
                        rit = i
                        bot = k
            if mx > ans:
                ans = mx
                a, b, c = lft, rit, bot
        return ans, a, b, c


    def center(self, l, t, w, h, imgW, imgH):
        l += w//2
        l -= imgW//2
        t += h//2
        t -= imgH//2
        return l, t


    def Process(self):
        totWidth = int(self.prs.slide_width.mm)+1
        totHeight = int(self.prs.slide_height.mm)+1
        for slide in range(1, len(self.prs.slides)):
            shapes = self.prs.slides[slide].shapes

            boxes = []
            for shape in self.prs.slides[slide].shapes:
                if hasattr(shape, "text"):
                    boxes.append((shape.left.mm, shape.top.mm, shape.width.mm, shape.height.mm))
            if not boxes:
                continue

            grid = [[1 for i in range(totWidth)] for j in range(totHeight)]
            for x, y, wi, hi in boxes:
                for i in range(int(y), int(y+hi)+1):
                    for j in range(int(x), int(x+wi)+1):
                        grid[i][j] = 0

            area, left, right, down = self.maximalRectangle(grid)
            width = right-left-self.margin-self.margin
            left += self.margin+1
            height = (area//width)-self.margin-self.margin
            top = down-height+1
            imageName = ""
            for filename in os.listdir(f'{self.image_dir}/Slide {slide+1}'):
                if filename.startswith("0"):
                    imageName = f'{self.image_dir}/Slide {slide+1}/{filename}'
            if not imageName:
                continue

            if height * 4 < width:
                picture = shapes.add_picture(imageName, 0, 0, width=Mm(totWidth))
                ratio = totHeight/picture.height.mm
                if ratio < 1:
                    picture.width = int(picture.width*ratio)
                    picture.height = int(picture.height*ratio)
                newL, newT = self.center(0, 0, totWidth, totHeight, picture.width.mm, picture.height.mm)
                picture.left = int(Mm(newL))
                picture.top = int(Mm(newT))
            else:
                picture = shapes.add_picture(imageName, Mm(left), Mm(top), width=Mm(width))
                ratio = height/picture.height.mm
                if ratio < 1:
                    picture.width = int(picture.width*ratio)
                    picture.height = int(picture.height*ratio)
                a, b = self.center(left, top, width, height, picture.width.mm, picture.height.mm)
                picture.left = int(Mm(a))
                picture.top = int(Mm(b))

        self.prs.save(f'{self.image_dir}/new-{self.filename}')


if __name__ == "__main__":
    prs_w = PresentationCreator("sample.pptx", "Sample Slides")
    prs_w.ReadPresentation()
    prs_w.Process()
