from pptx import Presentation
from pprint import pprint
from keyword_analyzer import KeywordAnalyzer, KeywordAnalyzer2, KeywordAnalyzer3
from webscraper import ImageScraper
import time
import asyncio

class PresentationReader:
    def __init__(self, presentation, download_dir="Slide Images", image_quality="thumb", complexity=2, debug=False):
        self.filename = presentation
        self.prs = Presentation(presentation)
        self.quality = image_quality
        self.download_dir = download_dir
        self.kw_analyzer = KeywordAnalyzer(max_ngram_size=complexity, keyword_count=complexity+1) #pass params
        self.scraper = ImageScraper(per_page=5, quality=self.quality, download_dir=self.download_dir) #pass params
        self.queries = [[] for _ in self.prs.slides]
        self.texts = ["" for _ in self.prs.slides]
        self.debug = debug

    def read_slides(self):
        for i, slide in enumerate(self.prs.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += (shape.text+"\n")
            self.texts[i] = text.strip()
        return self.texts

    def extract_keywords(self):
        for i, text in enumerate(self.texts):
            keywords = self.kw_analyzer.analyze_text(text)
            self.queries[i] = keywords        
        return self.queries

    def set_image_quality(self, quality):
        if quality in ["small", "regular", "thumb"]:
            self.quality = quality
            return True
        return False

    def set_download_dir(self, download_dir):
        self.download_dir = download_dir
    
    def find_images(self):
        self.scraper = ImageScraper(per_page=5, quality=self.quality, download_dir=self.download_dir)
        self.read_slides()
        self.extract_keywords()
        
        # Debug Output
        if self.debug:
            print(self.filename)
            print("----------------------")
            for i, slide in enumerate(self.prs.slides):
                # Debug Output
                print("Slide #", i)
                print("----------------------")
                print("Text: ")
                print(self.texts[i])
                print("Keywords: ")
                print(self.queries[i])

        start = time.perf_counter()
        self.scraper.download_images(self.queries)
        elapsed = time.perf_counter() - start
        print("ELAPSED TIME: ", elapsed)
            

if __name__ == "__main__":
    prs_r = PresentationReader("cheetah.pptx", complexity=1, debug=True)
    prs_r.set_download_dir("Cheetah Slides")
    prs_r.set_image_quality("thumb")
    prs_r.find_images()

    # prs = Presentation("weird_sample.pptx")
    # kw_analyzer = KeywordAnalyzer(max_ngram_size=1, keyword_count=3)
    # # kw_analyzer = KeywordAnalyzer2()
    # # kw_analyzer = KeywordAnalyzer3()
    # scraper = ImageScraper(per_page=5, quality="thumb", download_dir="Weird Images")
    # queries = [[] for _ in prs.slides]

    # # Opening powerpoint and extracting keywords
    # print("cheeto")
    # print("----------------------")

    # for i, slide in enumerate(prs.slides):
    #     print("Slide #", i)
    #     print("----------------------")

    #     text = ""
    #     for shape in slide.shapes:
    #         if hasattr(shape, "text"):
    #             # pprint(dir(shape))
    #             text += (shape.text+"\n")
    #     keywords = kw_analyzer.analyze_text(text)

    #     print("Text: ")
    #     print(text.strip())
    #     print("Keywords: ")
    #     print(keywords)

    #     queries[i] = keywords
    # print(queries)

    # # Getting images for keywords
    # scraper.download_images(queries)




